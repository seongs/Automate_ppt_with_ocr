import os
import re
from PIL import Image
import pytesseract
from pptx import Presentation
from pptx.util import Inches
import cv2


# 시리얼 번호 정규표현식
target_pattern = r"TH64C\s?\d+"

# PPT 파일 생성
presentation = Presentation()


# 전처리된 이미지로 OCR 수행

def extract_serial_number(image_path):
    """
    이미지에서 시리얼 번호를 추출
    """
    try:
        image = cv2.imread(image_path)
        gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY) # 이미지를 흑백으로 변환
        text = pytesseract.image_to_string(gray_image)
        print(text)
        match = re.findall(target_pattern, text)
        print(match)
        if match:
            return match[0]  # 시리얼 번호 반환
    except Exception as e:
        print(f"Error processing {image_path}: {e}")
    return None

def add_image_and_serial_to_ppt(slide, image_path, serial_number):
    """
    슬라이드에 이미지를 추가하고 시리얼 번호를 텍스트 삽입
    """
    # 이미지 추가
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(6))
    # 텍스트 추가
    left = Inches(1)
    top = Inches(5)
    textbox = slide.shapes.add_textbox(left, top, width=Inches(6), height=Inches(1))
    textbox.text = f"Serial Number: {serial_number}"

#이미지 디렉토리 설정
image_directory = "C:\\Users\\ksy\\Documents\\OCR\\imges"
images = [os.path.join(image_directory, f) for f in os.listdir(image_directory) if f.lower().endswith(('png', 'jpg', 'jpeg'))]

# 이미지 처리
for image_path in images:
    print(f"Image path: {image_path}")
    serial_number = extract_serial_number(image_path)
    if serial_number:
        # 슬라이드 추가
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # 빈 슬라이드
        add_image_and_serial_to_ppt(slide, image_path, serial_number)
        print(f"Added {image_path} with Serial Number: {serial_number}")

# 결과 PPT 저장
output_pptx = "Filtered_Images_With_Serial_Numbers.pptx"
presentation.save(output_pptx)
print(f"PPT saved as {output_pptx}") 